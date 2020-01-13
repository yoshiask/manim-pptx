import os
import platform
import glob
import subprocess
from manimlib.constants import *
import manimlib.addon_helper
from manimlib.scene.scene_file_writer import *

from pptx import Presentation

class Addon:
    PACKAGE_NAME = "pptx_export"

    CUR_DIR = os.path.join(ADDON_DIR, PACKAGE_NAME)
    TEMPLATE_DIR = os.path.join(CUR_DIR, "templates")
    TEMPLATE_PPTX = os.path.join(TEMPLATE_DIR, "template.pptx")
    EXAMPLE_PPTX = os.path.join(TEMPLATE_DIR, "powerpoint.pptx")
    TEMPORARY_DIR = os.path.join(CUR_DIR, "temporary")
    LOG_DIR = os.path.join(CUR_DIR, 'pptx_export.log')
        
    def addon_info(self):
        # Information about the addon
        return {
            "author": "Joshua \"Yoshi\" Askharoun",
            "name": Addon.PACKAGE_NAME,
            "version" : "1.0.1.0",
            "desc": "Adds --save_to_pptx, which generates a PowerPoint with a slide for each animation"
        }
    
    def cli_args(self):
        # Command-line flags added by the addon
        return [
            {
                "flag": "--save_to_pptx",
                "action": "store_true",
                "help": f"[{Addon.PACKAGE_NAME}] Render the animations to a PowerPoint presentation"
            },
            {
                "flag": "--anti_dupli_pptx",
                "action": "store_true",
                "help": f"[{Addon.PACKAGE_NAME}] When exporting to PowerPoint, only use every other movie part"
            },
            {
                "flag": "--join_scenes_pptx",
                "action": "store_false",
                "help": f"[{Addon.PACKAGE_NAME}] Combines the scenes into a single PowerPoint presentation"
            }
        ]
    
    def set_config(self, config):
        # Fires after the cli arguments are parsed in __init__.py for addons to have access to the current config
        self.config = config

    def loaded(self):
        # Fires when the addon is first initialized by manim
        #self.verbose("{0} addon loaded successfully".format(self.addon_info()['name']))
        return True

    def on_rendered(self, scene_classes):
        # Fires when a video is finished rendering
        args = self.config["all_args"]
        if args.save_to_pptx:
            if args.join_scenes_pptx:
                for scene_class in scene_classes:
                    name = scene_class.__name__
                    self.create_ppt([name], name)
            else:
                self.create_ppt([scene_class.__name__ for scene_class in scene_classes], "Scenes")

    def create_ppt(self, scene_names, ppt_name):
        SLD_BLANK = 6
        if not os.path.exists(Addon.TEMPORARY_DIR):
            os.makedirs(Addon.TEMPORARY_DIR)
        if os.path.exists(Addon.LOG_DIR):
            os.remove(Addon.LOG_DIR)
        prs = Presentation(Addon.TEMPLATE_PPTX)
        
        # Figure out where the movie parts are saved
        PART_DIRS = [os.path.join(os.path.dirname(manimlib.addon_helper.movie_paths[0]), "partial_movie_files", scene_name) for scene_name in scene_names]
        self.log_line(f"PART_DIRS = {PART_DIRS}")
        parts = [sorted(glob.glob(os.path.join(part_dir, "*.mp4"))) for part_dir in PART_DIRS]
        read_parts = PART_DIRS

        # Go through each video part and copy it over to the temporary directory.
        # If anti-duplicate is on, every other part is combined with the one before it and then copied over to the temp directory
        if self.config["all_args"].anti_dupli_pptx:
            self.log_line("Anti-duplication manually enabled")
            read_parts = [Addon.TEMPORARY_DIR]
            for file in parts:
                i = int(self.get_name(file))
                if not i % 2 == 0:
                    self.log_line("Merging parts {} and {}...".format(str(i), str(i-1)))
                    merged_clip = self.merge_videos(parts[i-1], file, os.path.join(Addon.TEMPORARY_DIR, str(i-1).zfill(5) + ".mp4"))
                    self.log_line("Merged to " + merged_clip)

        save_dir = os.path.join(os.path.dirname(manimlib.addon_helper.movie_paths[0]), ppt_name + ".pptx")
        slide_layout = prs.slide_layouts[SLD_BLANK]
        for part_dir in read_parts:
            for file in sorted(glob.glob(os.path.join(part_dir, "*.mp4"))):
                self.log_line(f"Using file at {file}")
                # Load the example presentation and its timing element
                prs_ex = Presentation(Addon.EXAMPLE_PPTX)
                timing_ex = prs_ex.slides[0].element[2]
                self.log_line(f"\tGrabbed timing element, timing_ex = {timing_ex}")
                # Create a new slide
                slide = prs.slides.add_slide(slide_layout)
                # Generate video thumbnail
                thumb_file = os.path.join(Addon.TEMPORARY_DIR, self.get_name(file) + ".png")
                self.log_line("\tGenerating video thumbnail...")
                self.get_video_thumb(file, thumb_file)
                self.log_line(f"\tThumbnail saved at {thumb_file}")

                # Add the video to the slide
                clip = slide.shapes.add_movie(file, 0, 0, prs.slide_width, prs.slide_height, mime_type='video/mp4', poster_frame_image=thumb_file)
                self.log_line("\tAdded clip to slide")

                # Play the clip in fullscreen when the slide starts
                ## Get the id of the movie object we just added
                id = clip.element[0][0].attrib.get("id")
                self.log_line(f"\tClip id = {id}")
                ## Make a copy of the timing element from the manually created PPTX,
                ## then change every spid to the clip id
                self.log_line("\tUsing timing_ex as template...")
                timing = timing_ex
                timing[0][0][0][0][0][0][0][0][0][1][0][0][1][0][0][1][0][0][1][0].attrib["spid"] = id
                timing[0][0][0][0][1][0][1][0].attrib["spid"] = id
                timing[0][0][0][0][2][0][0][0][0][0].attrib["spid"] = id
                timing[0][0][0][0][2][0][2][0][0][1][0][0][1][0][0][1][0][0][1][0].attrib["spid"] = id
                timing[0][0][0][0][2][1][0][0][0].attrib["spid"] = id
                slide.element[2] = timing
                self.log_line(f"\tAdded timing to slide, timing = {timing}")
                prs.save(save_dir)
                self.log_line(f"\tPPTX saved to {save_dir}")
        
        self.log_line(f"Final presentation saved to {save_dir}")
        print(f"\nPresentation ready at {save_dir}")
        if self.config["all_args"].preview:
            self.open_file(save_dir)

    def get_name(self, filename):
        pre, ext = os.path.splitext(filename)
        return pre.split(os.sep)[-1]

    def get_video_thumb(self, filename, imgname):
        command = [
                FFMPEG_BIN,
                '-y',  # overwrite output file if it exists
                '-loglevel', 'error',
                '-i', filename,
                '-vframes', '1',  # take only one frame
                imgname,
            ]
        subprocess.run(command, stdout=subprocess.PIPE)
        return imgname

    def get_frame_count(self, filename):
        command = [
            FFPROBE_BIN,
            '-v', 'error',
            '-count_frames',
            '-select_streams', 'v:0',
            '-show_entries', 'stream=nb_read_frames',
            '-of', 'default=nokey=1:noprint_wrappers=1',
            filename
        ]
        result = subprocess.run(command, stdout=subprocess.PIPE).stdout.decode('utf-8')
        return int(result)

    # TODO: Make this work
    def get_middle_video_frame(self, filename, imgname):
        command = [
            FFMPEG_BIN,
            '-y',  # overwrite output file if it exists
            '-loglevel', 'error',
            '-i', filename,
            '-vframes', '1',  # take only one frame
            imgname,
        ]
        subprocess.run(command, stdout=subprocess.PIPE)
        return imgname

    def merge_videos(self, clip1, clip2, output):
        vid_list = os.path.join(Addon.TEMPORARY_DIR, "cliplist.txt") 
        with open(vid_list, 'w') as file:
            file.write("file '{}'\n".format(clip1))
            file.write("file '{}'".format(clip2))
        commands = [
            FFMPEG_BIN,
            '-y',  # overwrite output file if it exists
            '-f', 'concat',
            '-safe', '0',
            '-i', vid_list,
            '-loglevel', 'error',
            '-c', 'copy', output
        ]
        subprocess.run(commands, stdout=subprocess.PIPE)
        return output

    def open_file(self, file_path):
        # Taken from open_file_if_needed()
        current_os = platform.system()
        if current_os == "Windows":
            os.startfile(file_path)
        else:
            commands = []
            if current_os == "Linux":
                commands.append("xdg-open")
            elif current_os.startswith("CYGWIN"):
                commands.append("cygstart")
            else:  # Assume macOS
                commands.append("open")

            if config["show_file_in_finder"]:
                commands.append("-R")

            commands.append(file_path)

            # commands.append("-g")
            FNULL = open(os.devnull, 'w')
            sp.call(commands, stdout=FNULL, stderr=sp.STDOUT)
            FNULL.close()

    def log_line(self, text):
        self.log_text(text.__str__() + "\n")

    def log_text(self, text):
        with open(Addon.LOG_DIR, 'a') as the_file:
            the_file.write(text.__str__())

    def __str__(self):
        return self.PACKAGE_NAME
    